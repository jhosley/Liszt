#!/usr/bin/env python3
"""
Seed the scenario library from the existing PowerPoint deck.

This is a ONE-TIME migration tool. It exists so the analytical work already
invested in the deck is not stranded when the repo becomes the system of record.
Run it once, review every generated record by hand, then delete your reliance on it.

    python tools/import_from_deck.py path/to/AIObservabilityAnalysis_Scenarios.pptx

Writes scenarios/NNN-slug.yaml for every scenario pair found.

Parsing strategy: pattern and geometry, never fixed shape indices. Shape indices
shift with step count, so index-based parsing silently mis-assigns fields.
"""
from __future__ import annotations

import argparse
import pathlib
import re
import sys
import unicodedata

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    sys.exit("python-pptx is required:  pip install python-pptx pyyaml")

try:
    import yaml
except ImportError:
    sys.exit("pyyaml is required:  pip install python-pptx pyyaml")


BASELINE = "2026.07"

TITLE_RE = re.compile(r"^Scenario\s+(\d+)\s*[·\-]\s*(.+)$")
STEP_RE = re.compile(r"^(\d)\s+\[(.+?)\]\s+(.+)$", re.S)

LAYER_MAP = {
    "L1 · Data": "L1 · Data",
    "L2 · Model": "L2 · Model",
    "L3 · Orchestration & Agent": "L3 · Orchestration & Agent",
    "L4 · Application": "L4 · Application",
}
EVIDENCE_MAP = {
    "seen in the wild": "seen-in-the-wild",
    "seen in research": "seen-in-research",
    "doomsday": "doomsday",
}
# Deck badge text -> current priority label. Old decks carry NEXT and LATER;
# re-importing one must produce the current labels NEAR-TERM and BACKLOG.
PRIORITY_MAP = {
    "NOW": "NOW",
    "NEXT": "NEAR-TERM",
    "LATER": "BACKLOG",
    "NEAR-TERM": "NEAR-TERM",
    "BACKLOG": "BACKLOG",
}
# Every badge text a deck may carry, old and new, for shape matching.
PRIORITY_BADGES = tuple(PRIORITY_MAP)


def norm(s: str) -> str:
    """Normalize whitespace and unicode without destroying the deck's · and → glyphs."""
    s = unicodedata.normalize("NFC", s or "")
    s = s.replace(" ", " ").replace("-", "-")
    return re.sub(r"[ \t]+", " ", s).strip()


def slugify(title: str) -> str:
    s = unicodedata.normalize("NFKD", title)
    s = s.replace("→", " to ").replace("&", " and ").replace("·", " ")
    s = re.sub(r"[^a-zA-Z0-9]+", "-", s).strip("-").lower()
    s = re.sub(r"-+", "-", s)
    return s[:60].strip("-")


def shape_text(sh) -> str:
    return norm(sh.text_frame.text) if sh.has_text_frame else ""


def texts(slide):
    """(shape, text, left, top) for every text-bearing shape, in reading order."""
    out = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = shape_text(sh)
        if t:
            out.append((sh, t, Emu(sh.left).inches, Emu(sh.top).inches))
    return sorted(out, key=lambda r: (round(r[3], 2), round(r[2], 2)))


# -- index slide --------------------------------------------------------------

def parse_index(prs):
    """Scenario number -> classification, from the index table."""
    for slide in prs.slides:
        for sh in slide.shapes:
            if not sh.has_table:
                continue
            tbl = sh.table
            hdr = [norm(tbl.cell(0, c).text).lower() for c in range(len(tbl.columns))]
            if "scenario" not in hdr or "priority" not in hdr:
                continue
            col = {name: i for i, name in enumerate(hdr)}
            index = {}
            for r in range(1, len(tbl.rows)):
                num = norm(tbl.cell(r, col["#"]).text)
                if not num.isdigit():
                    continue
                comp = norm(tbl.cell(r, col["primary layer component"]).text)
                layer = norm(tbl.cell(r, col["ai infrastructure layer"]).text)
                ev = norm(tbl.cell(r, col["evidence"]).text).lower()
                pri = norm(tbl.cell(r, col["priority"]).text).upper()
                index[num.zfill(3)] = {
                    "primary_layer_component": comp,
                    "ai_infrastructure_layer": LAYER_MAP.get(layer, layer),
                    "evidence": EVIDENCE_MAP.get(ev, ev.replace(" ", "-")),
                    "priority": PRIORITY_MAP.get(pri, "BACKLOG"),
                }
            if index:
                return index
    return {}


# -- attack-path slide --------------------------------------------------------

def parse_attack_path(slide):
    rec, steps = {}, []
    items = texts(slide)

    for _sh, t, _l, _top in items:
        m = TITLE_RE.match(t.split("\n")[0])
        if m and "telemetry" not in m.group(2).lower():
            rec["_num"] = m.group(1).zfill(3)
            rec["title"] = norm(m.group(2))
            break

    for _sh, t, left, _top in items:
        if left > 8.0:          # right-hand panel, not a step box
            continue
        m = STEP_RE.match(t)
        if m:
            steps.append({
                "step": int(m.group(1)),
                "layer": norm(m.group(2)),
                "text": norm(m.group(3).replace("\n", " ")),
            })
    steps.sort(key=lambda s: s["step"])
    rec["attack_path"] = steps

    # right-hand panel: In plain terms + How it maps
    for sh, t, _l, _top in items:
        if not t.startswith("In plain terms"):
            continue
        paras = [norm(p.text) for p in sh.text_frame.paragraphs]
        paras = [p for p in paras if p]
        if len(paras) > 1:
            rec["one_liner"] = paras[1]
        rec["_map_lines"] = paras[paras.index("How it maps") + 1:] if "How it maps" in paras else []
        break

    # bottom-left: priority rationale.  bottom-right: scaled_up.
    for sh, t, left, top in items:
        if top < 5.5:
            continue
        if left < 6.5 and t not in PRIORITY_BADGES and not t.startswith("Priority"):
            paras = [norm(p.text) for p in sh.text_frame.paragraphs if norm(p.text)]
            if paras and not paras[0].startswith("If this scaled up"):
                rec["priority_rationale"] = paras
        elif left >= 6.5 and not t.startswith("If this scaled up") and len(t) > 60:
            rec["scaled_up"] = t.replace("\n", " ")
    return rec


def parse_mapping(lines):
    """The three or four 'How it maps' bullets -> structured framework IDs."""
    out = {
        "baseline": BASELINE,
        "attack": [], "atlas": [], "owasp_llm": [], "owasp_agentic": [],
        "mapping_confidence": "editorial",
    }
    notes = []
    for ln in lines:
        low = ln.lower()
        out["attack"] += re.findall(r"\bT\d{4}(?:\.\d{3})?\b", ln)
        out["atlas"] += re.findall(r"\bAML\.(?:TA|T|M|CS)\d{4}(?:\.\d{3})?\b", ln)
        for n in re.findall(r"\bLLM(\d{2})\b", ln):
            out["owasp_llm"].append(f"LLM{n}:2025")
        for n in re.findall(r"\bASI(\d{2})\b", ln):
            out["owasp_agentic"].append(f"ASI{n}:2026")
        if low.startswith("atlas:") and not out["atlas"]:
            notes.append(ln)          # ATLAS named in prose, no ID, flag for the analyst
        if low.startswith("layers"):
            notes.append(ln)
    for k in ("attack", "atlas", "owasp_llm", "owasp_agentic"):
        out[k] = list(dict.fromkeys(out[k]))
    if notes:
        out["mapping_notes"] = (
            "IMPORTED FROM DECK, mapping was prose, not IDs. An analyst must resolve these "
            "to baseline " + BASELINE + " identifiers before this record is published: "
            + " | ".join(notes)
        )
    return out


# -- telemetry slide ----------------------------------------------------------

def parse_telemetry(slide):
    rec = {"telemetry": [], "commentary": {}, "_refs": []}
    items = texts(slide)

    for sh in slide.shapes:
        if not sh.has_table:
            continue
        tbl = sh.table
        for r in range(1, len(tbl.rows)):
            num = norm(tbl.cell(r, 0).text)
            if not num.isdigit():
                continue
            rec["telemetry"].append({
                "step": int(num),
                "signal": norm(tbl.cell(r, 1).text),
                "emitted_at": norm(tbl.cell(r, 2).text),
                "coverage": norm(tbl.cell(r, 3).text),
                "detection_opportunity": norm(tbl.cell(r, 4).text),
            })
        break

    for sh, t, left, top in items:
        if t.startswith("What we can already see:"):
            paras = [norm(p.text) for p in sh.text_frame.paragraphs if norm(p.text)]
            for p in paras:
                for prefix, key in (("What we can already see:", "already_see"),
                                    ("Where we're blind:", "blind"),
                                    ("How we detect it:", "how_detect")):
                    if p.startswith(prefix):
                        rec["commentary"][key] = norm(p[len(prefix):])
        elif 4.8 < top < 6.9 and left < 5.0 and "\n" in t and not t.startswith("Full references"):
            rec["_refs"] = [norm(p.text) for p in sh.text_frame.paragraphs if norm(p.text)]
    return rec


# -- assembly -----------------------------------------------------------------

def mark_control_rows(rows, steps):
    """
    A telemetry row numbered past the last attack step is not a missing step -- it is a
    control or verification signal the deck deliberately appended (e.g. "model signature
    verification"). Tag it so the validator reads the intent instead of flagging a gap.
    """
    n = len(steps)
    for r in rows:
        if r.get("step", 0) > n:
            r["kind"] = "control"
    return rows


def build(prs):
    index = parse_index(prs)
    pending, records = None, []

    for slide in prs.slides:
        title = ""
        for _sh, t, _l, _top in texts(slide):
            m = TITLE_RE.match(t.split("\n")[0])
            if m:
                title = t.split("\n")[0]
                break
        if not title:
            continue

        if "telemetry" in title.lower():
            if pending is None:
                print(f"  ! telemetry slide with no preceding attack-path slide: {title}")
                continue
            pending.update(parse_telemetry(slide))
            records.append(pending)
            pending = None
        else:
            if pending is not None:
                print(f"  ! attack-path slide with no telemetry slide: {pending.get('title')}")
                records.append(pending)
            pending = parse_attack_path(slide)

    if pending is not None:
        records.append(pending)

    out = []
    for r in records:
        num = r.get("_num")
        if not num:
            continue
        cls = index.get(num, {})
        rec = {
            "schema_version": 1,
            "id": num,
            "slug": slugify(r.get("title", "")),
            "title": r.get("title", ""),
            "one_liner": r.get("one_liner", ""),
            "status": "draft",          # imported records are NEVER born published
            "classification": {
                "primary_layer_component": cls.get("primary_layer_component", "Model"),
                "ai_infrastructure_layer": cls.get("ai_infrastructure_layer", "L2 · Model"),
                "evidence": cls.get("evidence", "seen-in-research"),
                "priority": cls.get("priority", "BACKLOG"),
                "priority_rationale": r.get("priority_rationale", []),
            },
            "framework_mapping": parse_mapping(r.get("_map_lines", [])),
            "attack_path": r.get("attack_path", []),
            "telemetry": mark_control_rows(r.get("telemetry", []), r.get("attack_path", [])),
            "commentary": r.get("commentary", {}),
            "scaled_up": r.get("scaled_up", ""),
            "incidents": [slugify(x) for x in r.get("_refs", []) if len(x) > 3],
            "provenance": {
                "authored_by": "imported-from-deck",
                "created": "2026-07-31",
                "last_updated": "2026-07-31",
                "sources": [],
            },
            "notes": (
                "IMPORTED from the source deck by tools/import_from_deck.py. "
                "Status is 'draft' on purpose. Before publishing: resolve framework IDs to the "
                "pinned baseline, add DeTT&CT scores and evidence to every telemetry row, name a "
                "data-source owner for each row, and have a second analyst review against "
                "docs/02-quality-bar.md."
            ),
        }
        out.append(rec)
    return out


def main():
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("deck", type=pathlib.Path)
    ap.add_argument("-o", "--out", type=pathlib.Path, default=pathlib.Path("scenarios"))
    ap.add_argument("--dry-run", action="store_true")
    args = ap.parse_args()

    if not args.deck.exists():
        sys.exit(f"not found: {args.deck}")

    records = build(Presentation(str(args.deck)))
    if not records:
        sys.exit("no scenario slide pairs found, is this the right deck?")

    args.out.mkdir(parents=True, exist_ok=True)
    for rec in records:
        path = args.out / f"{rec['id']}-{rec['slug']}.yaml"
        steps, rows = len(rec["attack_path"]), len(rec["telemetry"])
        flag = "" if steps == rows else f"   <-- {steps} steps vs {rows} rows, check by hand"
        print(f"  {path.name}{flag}")
        if not args.dry_run:
            with path.open("w", encoding="utf-8") as fh:
                yaml.safe_dump(rec, fh, sort_keys=False, allow_unicode=True, width=100)

    print(f"\n{len(records)} scenario records"
          f"{' (dry run, nothing written)' if args.dry_run else f' written to {args.out}/'}")
    print("Every one is status: draft. Nothing is published until a human reviews it.")


if __name__ == "__main__":
    main()
