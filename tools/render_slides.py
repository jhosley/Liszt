#!/usr/bin/env python3
"""
Render scenario records into the deck. The YAML is the source of truth; the deck
is a build artifact you can throw away and regenerate.

    # rebuild the whole deck from the library
    python tools/render_slides.py --template deck.pptx --out build/deck.pptx

    # one scenario only, appended to a copy of the template
    python tools/render_slides.py --template deck.pptx --out build/s21.pptx --only 021

    # refresh the index and incidents slides without touching scenario slides
    python tools/render_slides.py --template deck.pptx --out build/deck.pptx --index-only

How it works: the template deck supplies the visual design. We clone an existing
scenario slide pair and replace its text run by run, which preserves every font,
color and shape property without us having to know what they are. Never build
these slides from scratch, you will not reproduce the styling.
"""
from __future__ import annotations

import argparse
import copy
import pathlib
import re
import shutil
import subprocess
import sys
import tempfile
import zipfile

try:
    import yaml
    from pptx import Presentation
    from pptx.util import Emu, Inches
    from pptx.dml.color import RGBColor
except ImportError:
    sys.exit("pip install pyyaml python-pptx")

ROOT = pathlib.Path(__file__).resolve().parent.parent
# Appending NEW slide pairs needs tools/vendor/add_slide.py, which does all the package
# bookkeeping a new slide requires. Rendering into slide pairs the template already has
# needs nothing external, that is the common case, and it still works when the helper
# is absent.
ADD_SLIDE = next((p for p in (
    ROOT / "tools" / "vendor" / "add_slide.py",
) if p.exists()), None)

EVIDENCE_LABEL = {
    "seen-in-the-wild": "Seen in the wild",
    "seen-in-research": "Seen in research",
    "doomsday": "DOOMSDAY",
}

# Priority badge texts a template pair may carry: the union of the old labels
# (NEXT, LATER) and the current ones (NEAR-TERM, BACKLOG), so a template built
# before the label change still matches. The record's own label is what gets
# written into the badge.
PRIORITY_BADGES = ("NOW", "NEXT", "LATER", "NEAR-TERM", "BACKLOG")


# -- text helpers -------------------------------------------------------------
# Assigning to text_frame.text collapses a paragraph to one unstyled run and
# destroys the template's formatting. Always write into existing runs instead.

def set_para(shape, idx, text):
    p = shape.text_frame.paragraphs[idx]
    if not p.runs:
        return
    p.runs[0].text = text
    for r in p.runs[1:]:
        r._r.getparent().remove(r._r)


def set_runs(shape, idx, *values):
    p = shape.text_frame.paragraphs[idx]
    for run, val in zip(p.runs, values):
        run.text = val


def clone_para(shape, src_idx):
    tf = shape.text_frame
    new = copy.deepcopy(tf.paragraphs[src_idx]._p)
    tf.paragraphs[src_idx]._p.getparent().append(new)
    return tf.paragraphs[-1]


def fit_paras(shape, values, template_idx=None):
    """Grow or shrink a text frame to exactly len(values) paragraphs, then fill it."""
    tf = shape.text_frame
    src = template_idx if template_idx is not None else len(tf.paragraphs) - 1
    while len(tf.paragraphs) < len(values):
        clone_para(shape, src)
    while len(tf.paragraphs) > len(values):
        last = tf.paragraphs[-1]._p
        last.getparent().remove(last)
    for i, v in enumerate(values):
        set_para(shape, i, v)


def cell_text(cell, text):
    p = cell.text_frame.paragraphs[0]
    if not p.runs:
        return
    p.runs[0].text = text
    for r in p.runs[1:]:
        r._r.getparent().remove(r._r)


def find(slide, pred):
    for sh in slide.shapes:
        if sh.has_text_frame and pred(sh.text_frame.text.strip()):
            return sh
    return None


def by_geometry(slide, left_lt=None, left_ge=None, top_ge=None, top_lt=None, pred=None):
    out = []
    for sh in slide.shapes:
        if not sh.has_text_frame or not sh.text_frame.text.strip():
            continue
        l, t = Emu(sh.left).inches, Emu(sh.top).inches
        if left_lt is not None and l >= left_lt:
            continue
        if left_ge is not None and l < left_ge:
            continue
        if top_ge is not None and t < top_ge:
            continue
        if top_lt is not None and t >= top_lt:
            continue
        if pred and not pred(sh.text_frame.text.strip()):
            continue
        out.append((sh, l, t))
    return sorted(out, key=lambda r: (round(r[2], 2), round(r[1], 2)))


# -- slide writers ------------------------------------------------------------

STEP_RE = re.compile(r"^\d\s+\[.+?\]\s+", re.S)


def write_attack_path(slide, rec, page):
    set_para(find(slide, lambda t: t.startswith("Scenario")), 0,
             f"Scenario {rec['id'].lstrip('0') or '0'} · {rec['title']}")

    boxes = [sh for sh, _l, _t in by_geometry(slide, left_lt=8.0,
                                              pred=lambda t: STEP_RE.match(t))]
    steps = rec["attack_path"]
    if len(boxes) < len(steps):
        raise SystemExit(
            f"scenario {rec['id']}: template pair has {len(boxes)} step boxes but the record has "
            f"{len(steps)} steps. Clone a template pair that has at least {len(steps)}.")

    for box, s in zip(boxes, steps):
        set_runs(box, 0, f"{s['step']}  ", f"[{s['layer']}]  ", s["text"])

    # surplus boxes and their connecting arrows are hidden, not deleted -- deleting
    # shapes from a cloned slide is where template damage happens
    for box in boxes[len(steps):]:
        set_runs(box, 0, "", "", "")
        for sh in slide.shapes:
            if sh is not box and abs(Emu(sh.top).inches - Emu(box.top).inches) < 0.02 \
                    and abs(Emu(sh.left).inches - Emu(box.left).inches) < 0.02:
                sh.fill.background()
                sh.line.fill.background()

    panel = find(slide, lambda t: t.startswith("In plain terms"))
    fm = rec["framework_mapping"]
    bullets = []
    if fm.get("atlas"):
        bullets.append("ATLAS: " + " · ".join(fm["atlas"]))
    if fm.get("attack"):
        bullets.append("ATT&CK: " + " → ".join(fm["attack"]))
    owasp = fm.get("owasp_llm", []) + fm.get("owasp_agentic", [])
    if owasp:
        bullets.append("OWASP " + " + ".join(o.split(":")[0] for o in owasp))
    bullets.append(f"Layer {rec['classification']['ai_infrastructure_layer']}")

    hdr = [i for i, p in enumerate(panel.text_frame.paragraphs)
           if p.text.strip() == "How it maps"]
    body_start = hdr[0] + 1 if hdr else 3
    fit_paras(panel, ["In plain terms", rec["one_liner"], "How it maps"] + bullets,
              template_idx=body_start)

    badge = find(slide, lambda t: t in PRIORITY_BADGES)
    if badge:
        set_para(badge, 0, rec["classification"]["priority"])

    for sh, _l, _t in by_geometry(slide, left_lt=6.5, top_ge=5.5,
                                  pred=lambda t: t not in PRIORITY_BADGES
                                  and not t.startswith("Priority")):
        fit_paras(sh, rec["classification"]["priority_rationale"])
        break

    for sh, _l, _t in by_geometry(slide, left_ge=6.5, top_ge=5.5,
                                  pred=lambda t: not t.startswith("If this scaled up")
                                  and len(t) > 40):
        set_para(sh, 0, rec.get("scaled_up", ""))
        break

    num = find(slide, lambda t: t.isdigit() and len(t) <= 3)
    if num:
        set_para(num, 0, str(page))


def write_telemetry(slide, rec, page, incidents):
    set_para(find(slide, lambda t: t.startswith("Scenario")), 0,
             f"Scenario {rec['id'].lstrip('0') or '0'} · telemetry & detection map")
    subs = by_geometry(slide, top_lt=1.4, pred=lambda t: not t.startswith("Scenario"))
    if subs:
        set_para(subs[-1][0], 0, rec["title"])

    tbl = next(sh.table for sh in slide.shapes if sh.has_table)
    rows = rec["telemetry"]

    # reference cells for each coverage color, harvested from the template
    palette = {}
    for ri in range(1, len(tbl.rows)):
        palette.setdefault(tbl.cell(ri, 3).text.strip(), copy.deepcopy(tbl.cell(ri, 3)._tc))

    while len(tbl.rows) - 1 < len(rows):
        tbl._tbl.append(copy.deepcopy(tbl._tbl.tr_lst[-2]))
    while len(tbl.rows) - 1 > len(rows):
        tbl._tbl.remove(tbl._tbl.tr_lst[-1])

    for ri, r in enumerate(rows, start=1):
        for ci, val in enumerate([str(r["step"]), r["signal"], r["emitted_at"],
                                  r["coverage"], r["detection_opportunity"]]):
            if ci == 3 and val in palette:
                tr = tbl.cell(ri, ci)._tc.getparent()
                tr.replace(tbl.cell(ri, ci)._tc, copy.deepcopy(palette[val]))
            cell_text(tbl.cell(ri, ci), val)
    for row in tbl.rows:
        row.height = Inches(0.4 if len(rows) <= 6 else 0.34)

    badge = find(slide, lambda t: t in EVIDENCE_LABEL.values())
    if badge:
        set_para(badge, 0, EVIDENCE_LABEL[rec["classification"]["evidence"]])

    refs = [incidents.get(s, {}).get("title", s) for s in rec.get("incidents", [])] or \
           ["No incident of this kind has been observed"]
    for sh, _l, _t in by_geometry(slide, left_lt=5.0, top_ge=4.8, top_lt=6.9,
                                  pred=lambda t: not t.startswith("Full references")
                                  and not t.startswith("Seen in the real world")
                                  and t not in EVIDENCE_LABEL.values()):
        fit_paras(sh, refs)
        break

    com = find(slide, lambda t: t.startswith("What we can already see:"))
    if com:
        c = rec.get("commentary", {})
        for i, key in enumerate(("already_see", "blind", "how_detect")):
            if i < len(com.text_frame.paragraphs):
                p = com.text_frame.paragraphs[i]
                if len(p.runs) >= 2:
                    p.runs[1].text = c.get(key, "") + ("\n" if i < 2 else "")

    num = find(slide, lambda t: t.isdigit() and len(t) <= 3)
    if num:
        set_para(num, 0, str(page))


# -- deck assembly ------------------------------------------------------------

def scenario_slide_pairs(prs):
    pairs, pending = [], None
    for i, slide in enumerate(prs.slides):
        title = next((sh.text_frame.text.strip() for sh in slide.shapes
                      if sh.has_text_frame and sh.text_frame.text.strip().startswith("Scenario")),
                     "")
        if not title:
            continue
        if "telemetry" in title.lower():
            if pending is not None:
                pairs.append((pending, i))
                pending = None
        else:
            pending = i
    return pairs


def clone_pair(work: pathlib.Path, src_a: str, src_b: str, after: str):
    if ADD_SLIDE is None:
        sys.exit(
            "This render needs to APPEND a new slide pair, which requires the helper at\n"
            "tools/vendor/add_slide.py (it performs the content-type, relationship and\n"
            "sldIdLst registration a new slide needs; copying the XML by hand corrupts\n"
            "the deck).\n\n"
            "Either:\n"
            "  - use a template deck that already contains a slide pair for every scenario\n"
            "    you are rendering (--only <ids> to render just those), or\n"
            "  - place add_slide.py at tools/vendor/add_slide.py.")
    for src in (src_a, src_b):
        out = subprocess.run([sys.executable, str(ADD_SLIDE), str(work), src, "--after", after],
                             capture_output=True, text=True)
        if out.returncode:
            sys.exit(f"add_slide failed: {out.stderr}")
        after = re.search(r"Created ppt/slides/(slide\d+\.xml)", out.stdout).group(1)
    return after


def main():
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--template", required=True, type=pathlib.Path)
    ap.add_argument("--out", required=True, type=pathlib.Path)
    ap.add_argument("--only", nargs="*", help="scenario ids, e.g. 021 022")
    ap.add_argument("--include-drafts", action="store_true",
                    help="by default only status: published records are rendered")
    args = ap.parse_args()

    incidents = {}
    for p in (ROOT / "incidents").glob("*.yaml"):
        incidents[p.stem] = yaml.safe_load(p.read_text()) or {}

    records = []
    for p in sorted((ROOT / "scenarios").glob("*.yaml")):
        if p.name.startswith("_"):        # _TEMPLATE.yaml is not a record
            continue
        rec = yaml.safe_load(p.read_text())
        if args.only and rec["id"] not in args.only:
            continue
        if rec.get("status") == "retired":
            continue
        if rec.get("status") != "published" and not args.include_drafts:
            continue
        records.append(rec)

    if not records:
        sys.exit("nothing to render. Records are only rendered when status is 'published' "
                 "-- pass --include-drafts to preview work in progress.")

    args.out.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy(args.template, args.out)

    prs = Presentation(str(args.out))
    pairs = scenario_slide_pairs(prs)
    if not pairs:
        sys.exit("template contains no scenario slide pair to clone")

    # widest existing pair becomes the clone source, so any record fits
    widest = max(pairs, key=lambda pr: len(by_geometry(prs.slides[pr[0]], left_lt=8.0,
                                                       pred=lambda t: STEP_RE.match(t))))
    src_a, src_b = f"slide{widest[0] + 1}.xml", f"slide{widest[1] + 1}.xml"
    last_scenario_slide = f"slide{pairs[-1][1] + 1}.xml"

    with tempfile.TemporaryDirectory() as td:
        work = pathlib.Path(td) / "work.pptx"
        shutil.copy(args.out, work)

        existing = {}
        for a, b in pairs:
            t = next(sh.text_frame.text for sh in prs.slides[a].shapes
                     if sh.has_text_frame and sh.text_frame.text.startswith("Scenario"))
            m = re.match(r"Scenario\s+(\d+)", t)
            if m:
                existing[m.group(1).zfill(3)] = (a, b)

        after = last_scenario_slide
        appended = []
        for rec in records:
            if rec["id"] not in existing:
                after = clone_pair(work, src_a, src_b, after)
                appended.append(rec["id"])

        prs = Presentation(str(work))
        pairs = scenario_slide_pairs(prs)
        by_id, unused = {}, []
        for a, b in pairs:
            t = next(sh.text_frame.text for sh in prs.slides[a].shapes
                     if sh.has_text_frame and sh.text_frame.text.startswith("Scenario"))
            m = re.match(r"Scenario\s+(\d+)", t)
            key = m.group(1).zfill(3) if m else None
            (by_id.setdefault(key, (a, b)) if key not in by_id else unused.append((a, b)))

        page = None
        for rec in records:
            a, b = by_id.get(rec["id"]) or unused.pop(0)
            if page is None:
                cur = next((sh.text_frame.text for sh in prs.slides[a].shapes
                            if sh.has_text_frame and sh.text_frame.text.strip().isdigit()), "1")
                page = int(cur)
            write_attack_path(prs.slides[a], rec, page)
            write_telemetry(prs.slides[b], rec, page + 1, incidents)
            page += 2
            print(f"  {rec['id']}  {rec['title']}")

        prs.save(str(args.out))

    print(f"\n{len(records)} scenario(s) rendered to {args.out}")
    if appended:
        print(f"appended new slide pairs for: {', '.join(appended)}")


if __name__ == "__main__":
    main()
